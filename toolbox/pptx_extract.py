"""Extract pptx slide-by-slide text, tables, notes; dump images to assets/."""
import sys, os, json, re
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu

src = Path(sys.argv[1])
out_dir = Path(sys.argv[2])
assets_dir = out_dir / "assets" / src.stem
assets_dir.mkdir(parents=True, exist_ok=True)

prs = Presentation(str(src))

def shape_text(shape):
    parts = []
    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            line = "".join(r.text for r in p.runs) or p.text
            if line.strip():
                parts.append(line.rstrip())
    if shape.has_table:
        for row in shape.table.rows:
            cells = [c.text_frame.text.strip().replace("\n", " / ") for c in row.cells]
            parts.append(" | ".join(cells))
    return parts

def walk(shape, lines, imgs, slide_idx):
    if shape.shape_type == 6:  # group
        for s in shape.shapes:
            walk(s, lines, imgs, slide_idx)
        return
    if shape.shape_type == 13 or getattr(shape, "image", None):
        try:
            img = shape.image
            ext = img.ext
            name = f"slide{slide_idx:02d}_{len(imgs)+1}.{ext}"
            (assets_dir / name).write_bytes(img.blob)
            imgs.append(name)
        except Exception:
            pass
    lines.extend(shape_text(shape))

md_lines = [f"# {src.stem}", ""]
md_lines.append(f"- Slides: {len(prs.slides)}")
md_lines.append(f"- Source: `{src.name}`")
md_lines.append(f"- Assets: `assets/{src.stem}/`")
md_lines.append("")

for i, slide in enumerate(prs.slides, 1):
    title = ""
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        title = slide.shapes.title.text_frame.text.strip().replace("\n", " ")
    md_lines.append(f"## Slide {i}" + (f" — {title}" if title else ""))
    md_lines.append("")
    lines, imgs = [], []
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        try:
            walk(shape, lines, imgs, i)
        except Exception as e:
            lines.append(f"<!-- shape err: {e} -->")
    for l in lines:
        md_lines.append(f"- {l}")
    if imgs:
        md_lines.append("")
        md_lines.append("**Images:**")
        for name in imgs:
            md_lines.append(f"- ![](assets/{src.stem}/{name})")
    notes = ""
    if slide.has_notes_slide:
        n = slide.notes_slide.notes_text_frame.text.strip()
        if n:
            notes = n
    if notes:
        md_lines.append("")
        md_lines.append("**Notes:**")
        for line in notes.splitlines():
            if line.strip():
                md_lines.append(f"> {line}")
    md_lines.append("")

out_file = out_dir / f"{src.stem}_extract.md"
out_file.write_text("\n".join(md_lines), encoding="utf-8")
print(f"Wrote {out_file}")
print(f"Images: {len(list(assets_dir.iterdir()))}")
